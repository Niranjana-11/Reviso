"""
pptx_converter.py
Converts a .pptx file to a .pdf file using python-pptx + reportlab.
Extracts all text from every slide and writes into a clean PDF.
"""

from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
import re


def _safe_text(text: str) -> str:
    """
    Escape special HTML characters that break ReportLab's paragraph parser.
    ReportLab uses XML-like parsing so <, >, & must be escaped.
    """
    text = text.replace("&", "&amp;")
    text = text.replace("<", "&lt;")
    text = text.replace(">", "&gt;")
    text = text.replace('"', "&quot;")
    text = text.replace("'", "&#39;")
    return text


def pptx_to_pdf(pptx_path: str, output_path: str) -> str:
    """
    Read all text from a PPTX and write it into a clean PDF.
    Returns the output PDF path.
    """
    prs = Presentation(pptx_path)

    # ── Extract text slide by slide ───────────────────────────────────────
    slides_content = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
        if texts:
            slides_content.append((i, texts))

    if not slides_content:
        raise ValueError("No readable text found in the PowerPoint file.")

    # ── Build PDF ─────────────────────────────────────────────────────────
    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm,
    )

    base = getSampleStyleSheet()

    slide_title_style = ParagraphStyle(
        "SlideTitle",
        parent=base["Normal"],
        fontSize=13,
        fontName="Helvetica-Bold",
        textColor=colors.HexColor("#1e3a5f"),
        spaceAfter=4,
    )
    body_style = ParagraphStyle(
        "SlideBody",
        parent=base["Normal"],
        fontSize=10.5,
        textColor=colors.HexColor("#374151"),
        leading=15,
        leftIndent=10,
    )
    rule_color = colors.HexColor("#e2e8f0")

    story = []

    for slide_num, texts in slides_content:
        # Slide number header
        story.append(
            Paragraph(f"Slide {slide_num}", slide_title_style)
        )

        for j, line in enumerate(texts):
            # ← KEY FIX: escape special characters before passing to ReportLab
            safe_line = _safe_text(line)
            style = slide_title_style if j == 0 else body_style
            try:
                story.append(Paragraph(safe_line, style))
            except Exception:
                # If paragraph still fails, use plain text fallback
                plain = re.sub(r'[^\x20-\x7E]', ' ', safe_line)
                story.append(Paragraph(plain, style))

        story.append(Spacer(1, 0.2*cm))
        story.append(HRFlowable(
            width="100%", thickness=0.5, color=rule_color
        ))
        story.append(Spacer(1, 0.3*cm))

    doc.build(story)
    return output_path